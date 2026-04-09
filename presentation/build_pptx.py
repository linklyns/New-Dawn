#!/usr/bin/env python3
"""
Build the New Dawn INTEX presentation PowerPoint.
Section 3, Group 9 — LIGHT THEME
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (Light Theme) ────────────────────────────────────
SLATE_NAVY   = RGBColor(0x2D, 0x3A, 0x4A)
DARK_NAVY    = RGBColor(0x1A, 0x26, 0x33)
SKY_BLUE     = RGBColor(0x6A, 0xA8, 0xCB)   # slightly deeper for light bg contrast
SAGE_GREEN   = RGBColor(0x5E, 0x8C, 0x5E)   # deeper green for readability
GOLDEN_HONEY = RGBColor(0xE6, 0xA8, 0x17)   # deeper gold for light bg
CORAL_PINK   = RGBColor(0xFF, 0xE6, 0xE1)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF3, 0xF4, 0xF6)
CARD_BG      = RGBColor(0xF7, 0xF8, 0xFA)
MED_GRAY     = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT     = RGBColor(0x2D, 0x3A, 0x4A)   # main text = slate navy
BODY_TEXT     = RGBColor(0x55, 0x5F, 0x6D)   # secondary body text
LIGHT_BODY   = RGBColor(0x6B, 0x72, 0x80)   # lighter body
ACCENT_BAR   = SLATE_NAVY

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


# ── Helpers ─────────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape


def add_rounded_rect(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=DARK_TEXT, font_name="Calibri",
                    bullet_color=None, spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.line_spacing = Pt(int(font_size * spacing))
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.findall(qn("a:buNone"))
        for bn in buNone:
            pPr.remove(bn)
        from lxml import etree
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")
        if bullet_color:
            buClr = etree.SubElement(pPr, qn("a:buClr"))
            srgb = etree.SubElement(buClr, qn("a:srgbClr"))
            srgb.set("val", "%02X%02X%02X" % (bullet_color[0], bullet_color[1], bullet_color[2]))
    return txBox


def add_numbered_list(slide, items, left, top, width, height,
                      font_size=14, color=DARK_TEXT, font_name="Calibri",
                      num_color=None, spacing=1.4):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        from pptx.oxml.ns import qn
        from lxml import etree
        p.clear()
        r1 = p.add_run()
        r1.text = f"{i+1}. "
        r1.font.size = Pt(font_size)
        r1.font.color.rgb = num_color or color
        r1.font.name = font_name
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = color
        r2.font.name = font_name
        r2.font.bold = False
        p.line_spacing = Pt(int(font_size * spacing))
    return txBox


def add_accent_line(slide, left, top, width, color=SLATE_NAVY, height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def slide_header(slide, label, title, accent_color=SLATE_NAVY, title_color=DARK_TEXT):
    add_bg(slide, WHITE)
    add_shape(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BAR)
    add_text(slide, label, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             font_size=14, color=accent_color, bold=True)
    add_text(slide, title, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
             font_size=38, color=title_color, bold=True, font_name="Calibri Light")
    add_accent_line(slide, Inches(0.8), Inches(1.65), Inches(2), accent_color, Pt(3))


def add_slide_number(slide, num, total):
    add_text(slide, f"{num} / {total}", Inches(12.2), Inches(7.05),
             Inches(1), Inches(0.4), font_size=10, color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 10


# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BAR)

add_text(slide, "NEW DAWN", Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         font_size=64, color=SLATE_NAVY, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")

add_accent_line(slide, Inches(5.5), Inches(3.1), Inches(2.3), GOLDEN_HONEY, Pt(3))

add_text(slide, "A Path to Healing and Hope", Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
         font_size=28, color=SKY_BLUE, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

add_text(slide, "Section 3  |  Group 9", Inches(1.5), Inches(5.4), Inches(10), Inches(0.5),
         font_size=16, color=SLATE_NAVY, alignment=PP_ALIGN.CENTER, bold=True)

add_text(slide, "Zack Hada   |   Lincoln Lyons   |   [Team Members]",
         Inches(1.5), Inches(5.85), Inches(10), Inches(0.5),
         font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — THE MISSION (sensitive intro)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BAR)

add_text(slide,
         "YOUR MISSION", Inches(1.5), Inches(1.0), Inches(10.3), Inches(0.5),
         font_size=14, color=SLATE_NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide,
         "You're building an organization to provide safe homes for girls\nwho are survivors of sexual abuse and trafficking in the Philippines.",
         Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
         font_size=26, color=DARK_TEXT, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light", line_spacing=1.5)

add_text(slide,
         "Caring. Healing. Teaching.",
         Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.7),
         font_size=22, color=SAGE_GREEN, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light", bold=True)

add_accent_line(slide, Inches(5.8), Inches(3.8), Inches(1.7), GOLDEN_HONEY, Pt(2))

add_text(slide,
         "Shelter, counseling, education, and reintegration\nfor the most vulnerable -- with limited staff, limited budget,\nand data that could be saving lives if you had the tools to use it.",
         Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.5),
         font_size=20, color=BODY_TEXT, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light", line_spacing=1.5)

add_text(slide,
         "You asked us to build those tools.",
         Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
         font_size=22, color=SLATE_NAVY, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light", bold=True)

add_text(slide,
         "This is what we built.",
         Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.6),
         font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")

add_slide_number(slide, 2, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — THE CHALLENGES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "WHAT YOU SHARED WITH US", "Three Things That Keep You Up at Night")

card_y = Inches(2.2)
card_h = Inches(4.5)
card_w = Inches(3.6)
gap = Inches(0.4)
sx = Inches(0.8)

# Card 1
c1 = add_rounded_rect(slide, sx, card_y, card_w, card_h, fill=CARD_BG)
add_shape(slide, sx, card_y, card_w, Inches(0.06), fill=SAGE_GREEN)
add_text(slide, "\"We worry about girls\nfalling through the cracks.\"",
         sx + Inches(0.3), card_y + Inches(0.4),
         Inches(3), Inches(1.0), font_size=17, color=SAGE_GREEN, bold=True,
         font_name="Calibri Light")
add_text(slide, "With limited staff across multiple safehouses, you can't always see which girls are progressing and which are struggling. There's no way to predict when someone is ready for reintegration -- or at risk of regression.",
         sx + Inches(0.3), card_y + Inches(1.6),
         Inches(3), Inches(2.6), font_size=14, color=BODY_TEXT)

# Card 2
c2x = sx + card_w + gap
c2 = add_rounded_rect(slide, c2x, card_y, card_w, card_h, fill=CARD_BG)
add_shape(slide, c2x, card_y, card_w, Inches(0.06), fill=SKY_BLUE)
add_text(slide, "\"We don't know what\nto post or when.\"",
         c2x + Inches(0.3), card_y + Inches(0.4),
         Inches(3), Inches(1.0), font_size=17, color=SKY_BLUE, bold=True,
         font_name="Calibri Light")
add_text(slide, "Social media is your primary channel for reaching donors, but you've told us you aren't experienced with it. You post sporadically and can't tell what actually drives donations versus just generating likes.",
         c2x + Inches(0.3), card_y + Inches(1.6),
         Inches(3), Inches(2.6), font_size=14, color=BODY_TEXT)

# Card 3
c3x = c2x + card_w + gap
c3 = add_rounded_rect(slide, c3x, card_y, card_w, card_h, fill=CARD_BG)
add_shape(slide, c3x, card_y, card_w, Inches(0.06), fill=GOLDEN_HONEY)
add_text(slide, "\"We lose donors and\ndon't understand why.\"",
         c3x + Inches(0.3), card_y + Inches(0.4),
         Inches(3), Inches(1.0), font_size=17, color=GOLDEN_HONEY, bold=True,
         font_name="Calibri Light")
add_text(slide, "Your organization depends entirely on donations. You run fundraising campaigns but aren't sure which ones move the needle. You want to know which donors might give more, which are at risk of lapsing, and how to personalize outreach without a marketing team.",
         c3x + Inches(0.3), card_y + Inches(1.6),
         Inches(3), Inches(2.6), font_size=14, color=BODY_TEXT)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — OUR SOLUTION
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "OUR ANSWER", "One Platform, Three Solutions")

add_text(slide, "We built New Dawn -- a full-stack web application that combines case management,\ndonor intelligence, and AI-powered social media tools into one system.",
         Inches(0.8), Inches(2.0), Inches(11), Inches(0.8),
         font_size=18, color=BODY_TEXT)

sol_y = Inches(3.2)
sol_h = Inches(3.0)
sol_w = Inches(3.6)

# Solution 1
s1 = add_rounded_rect(slide, sx, sol_y, sol_w, sol_h, fill=CARD_BG)
add_shape(slide, sx, sol_y, sol_w, Inches(0.06), fill=SAGE_GREEN)
add_text(slide, "Caseload Intelligence", sx + Inches(0.3), sol_y + Inches(0.3),
         Inches(3), Inches(0.5), font_size=20, color=SAGE_GREEN, bold=True)
add_text(slide, "ML risk prediction flags at-risk residents before they regress. Causal analysis identifies which factors actually drive successful reintegration. Automated notifications alert staff to forgotten cases.",
         sx + Inches(0.3), sol_y + Inches(0.9),
         Inches(3), Inches(1.5), font_size=14, color=BODY_TEXT)
add_text(slide, "Risk Prediction + Reintegration Causal + Notifications",
         sx + Inches(0.3), sol_y + Inches(2.3),
         Inches(3), Inches(0.5), font_size=11, color=MED_GRAY)

# Solution 2
s2x = sx + sol_w + gap
s2 = add_rounded_rect(slide, s2x, sol_y, sol_w, sol_h, fill=CARD_BG)
add_shape(slide, s2x, sol_y, sol_w, Inches(0.06), fill=SKY_BLUE)
add_text(slide, "AI Social Editor", s2x + Inches(0.3), sol_y + Inches(0.3),
         Inches(3), Inches(0.5), font_size=20, color=SKY_BLUE, bold=True)
add_text(slide, "AI generates and refines social media posts via interactive chat. ML predicts donations, engagement, and reach in real time. Optimal posting times maximize fundraising. Save drafts, upload media, preview per platform.",
         s2x + Inches(0.3), sol_y + Inches(0.9),
         Inches(3), Inches(1.5), font_size=14, color=BODY_TEXT)
add_text(slide, "AI Generation + Referral Predictions + Best Posting Times",
         s2x + Inches(0.3), sol_y + Inches(2.3),
         Inches(3), Inches(0.5), font_size=11, color=MED_GRAY)

# Solution 3
s3x = s2x + sol_w + gap
s3 = add_rounded_rect(slide, s3x, sol_y, sol_w, sol_h, fill=CARD_BG)
add_shape(slide, s3x, sol_y, sol_w, Inches(0.06), fill=GOLDEN_HONEY)
add_text(slide, "Donor Intelligence", s3x + Inches(0.3), sol_y + Inches(0.3),
         Inches(3), Inches(0.5), font_size=20, color=GOLDEN_HONEY, bold=True)
add_text(slide, "ML predicts which donors are likely to give again within 6 months, with SHAP-driven reasons explaining each score. Automated alerts flag declining donors so staff can act before they lapse.",
         s3x + Inches(0.3), sol_y + Inches(0.9),
         Inches(3), Inches(1.5), font_size=14, color=BODY_TEXT)
add_text(slide, "Donor Likelihood Prediction + Decline Notifications",
         s3x + Inches(0.3), sol_y + Inches(2.3),
         Inches(3), Inches(0.5), font_size=11, color=MED_GRAY)

add_text(slide, "Let us show you.",
         Inches(1.5), Inches(6.6), Inches(10), Inches(0.5),
         font_size=20, color=SLATE_NAVY, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light", bold=True)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — CONSOLIDATED LIVE DEMO GUIDE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BAR)

add_text(slide, "LIVE DEMO", Inches(0.8), Inches(0.3), Inches(3), Inches(0.5),
         font_size=14, color=SLATE_NAVY, bold=True)
add_text(slide, "new-dawn-virid.vercel.app", Inches(8.5), Inches(0.3), Inches(4.5), Inches(0.5),
         font_size=14, color=SKY_BLUE, alignment=PP_ALIGN.RIGHT)

col_w = Inches(3.85)
col_gap = Inches(0.25)
col_top = Inches(0.9)
col_h = Inches(6.3)

# ── COLUMN 1: Public + Admin ──
c1x = Inches(0.5)
add_rounded_rect(slide, c1x, col_top, col_w, col_h, fill=CARD_BG)

add_text(slide, "PUBLIC EXPERIENCE", c1x + Inches(0.25), col_top + Inches(0.2),
         Inches(3.3), Inches(0.35), font_size=12, color=SLATE_NAVY, bold=True)

pub_items = [
    "Landing page -- mission, three pillars (Caring, Healing, Teaching), live stats",
    "Impact Dashboard -- charts: residents served, program outcomes, donation impact",
    "Donate page -- four tiers with impact messaging",
    "Dark/light mode toggle (stored in browser cookie)",
    "Privacy policy + GDPR cookie consent banner",
]
add_numbered_list(slide, pub_items, c1x + Inches(0.25), col_top + Inches(0.6),
                  Inches(3.4), Inches(2.6),
                  font_size=12, color=BODY_TEXT, num_color=SLATE_NAVY, spacing=1.6)

add_accent_line(slide, c1x + Inches(0.25), col_top + Inches(3.2), Inches(3.3),
                LIGHT_GRAY, Pt(1))

add_text(slide, "AUTH & ADMIN", c1x + Inches(0.25), col_top + Inches(3.4),
         Inches(3.3), Inches(0.35), font_size=12, color=SLATE_NAVY, bold=True)

admin_items = [
    "Login with Google OAuth; show MFA setup in profile",
    "Admin Dashboard -- residents, donations, OKR metric, reintegration insights",
    "Notification bell -- alerts for at-risk donors, forgotten cases, milestones",
    "Role-based access, delete confirmations, pagination on every list",
]
add_numbered_list(slide, admin_items, c1x + Inches(0.25), col_top + Inches(3.8),
                  Inches(3.4), Inches(2.3),
                  font_size=12, color=BODY_TEXT, num_color=SLATE_NAVY, spacing=1.6)

# ── COLUMN 2: Caseload + Donors ──
c2x = c1x + col_w + col_gap
add_rounded_rect(slide, c2x, col_top, col_w, col_h, fill=CARD_BG)

add_text(slide, "CASELOAD INTELLIGENCE", c2x + Inches(0.25), col_top + Inches(0.2),
         Inches(3.3), Inches(0.35), font_size=12, color=SAGE_GREEN, bold=True)

case_items = [
    "Residents list -- filters by status, safehouse, risk. Show ML risk arrows",
    "Click a resident -- full profile, demographics, case info, disability data",
    "ML Risk Assessment card -- predicted risk + confidence + top SHAP factors",
    "Reintegration Insights -- causal factors from logistic regression",
    "Process recordings, education, health, interventions, incidents (show 1-2)",
]
add_numbered_list(slide, case_items, c2x + Inches(0.25), col_top + Inches(0.6),
                  Inches(3.4), Inches(2.6),
                  font_size=12, color=BODY_TEXT, num_color=SAGE_GREEN, spacing=1.6)

add_accent_line(slide, c2x + Inches(0.25), col_top + Inches(3.2), Inches(3.3),
                LIGHT_GRAY, Pt(1))

add_text(slide, "DONOR INTELLIGENCE", c2x + Inches(0.25), col_top + Inches(3.4),
         Inches(3.3), Inches(0.35), font_size=12, color=GOLDEN_HONEY, bold=True)

donor_items = [
    "Supporters list with ML likelihood badges (High/Med/Low)",
    "SHAP reasons in plain language explain each prediction",
    "Sort by likelihood to prioritize at-risk donor outreach",
    "Donor detail -- history, lifetime value, frequency, avg gift",
    "Allocations -- where each dollar goes by safehouse/program",
]
add_numbered_list(slide, donor_items, c2x + Inches(0.25), col_top + Inches(3.8),
                  Inches(3.4), Inches(2.3),
                  font_size=12, color=BODY_TEXT, num_color=GOLDEN_HONEY, spacing=1.6)

# ── COLUMN 3: Social Editor + Reports ──
c3x = c2x + col_w + col_gap
add_rounded_rect(slide, c3x, col_top, col_w, col_h, fill=CARD_BG)

add_shape(slide, c3x + Inches(2.4), col_top + Inches(0.15), Inches(1.2), Inches(0.3),
          fill=GOLDEN_HONEY)
add_text(slide, "HIGHLIGHT", c3x + Inches(2.45), col_top + Inches(0.17),
         Inches(1.1), Inches(0.25), font_size=10, color=WHITE, bold=True)

add_text(slide, "AI SOCIAL EDITOR", c3x + Inches(0.25), col_top + Inches(0.2),
         Inches(2.2), Inches(0.35), font_size=12, color=SKY_BLUE, bold=True)

social_items = [
    "Fill out a brief (platform, topic, CTA, tone) then AI generates the post",
    "Refine via AI chat -- ask it to adjust tone, add hashtags, etc.",
    "Upload media and see platform-specific previews (IG Reel, TikTok, etc.)",
    "Predict Performance: 6 ML metrics update (referrals, value, engagement...)",
    "Suggest Times: ML-ranked optimal day/hour slots, click to schedule",
]
add_numbered_list(slide, social_items, c3x + Inches(0.25), col_top + Inches(0.6),
                  Inches(3.4), Inches(2.6),
                  font_size=12, color=BODY_TEXT, num_color=SKY_BLUE, spacing=1.6)

add_accent_line(slide, c3x + Inches(0.25), col_top + Inches(3.2), Inches(3.3),
                LIGHT_GRAY, Pt(1))

add_text(slide, "REPORTS & ANALYTICS", c3x + Inches(0.25), col_top + Inches(3.4),
         Inches(3.3), Inches(0.35), font_size=12, color=SLATE_NAVY, bold=True)

report_items = [
    "Donation trends (monthly totals + counts over time)",
    "Education progress & health trends by safehouse",
    "Reintegration success rates per safehouse",
    "Incident summary by type and severity",
    "Social media post performance and platform comparison",
]
add_numbered_list(slide, report_items, c3x + Inches(0.25), col_top + Inches(3.8),
                  Inches(3.4), Inches(2.3),
                  font_size=12, color=BODY_TEXT, num_color=SLATE_NAVY, spacing=1.6)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — ML PIPELINES OVERVIEW
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "UNDER THE HOOD", "5 Production ML Pipelines")

pipelines = [
    ("1", "Donor Likelihood", "Predictive", "Gradient Boosting", "F1: 0.98, AUC: 0.98", "Supporters list badges"),
    ("2", "Social Referrals", "Explanatory + Predictive", "GBM / Ridge / LightGBM", "6 target models", "Social Editor live cards"),
    ("3", "Best Posting Times", "Predictive", "Regression", "168 day/hr combos", "Social Editor scheduler"),
    ("4", "Reintegration Causal", "Explanatory", "Logistic (statsmodels)", "Odds ratios, p<0.05", "Resident detail insights"),
    ("5", "Risk Prediction", "Predictive", "Random Forest", "Accuracy: 65%", "Residents list + detail"),
]

header_y = Inches(2.0)
add_shape(slide, Inches(0.6), header_y, Inches(12.1), Inches(0.45), fill=SLATE_NAVY)
cols = [("#", 0.6, 0.4), ("Pipeline", 1.1, 2.2), ("Approach", 3.4, 2.2), ("Algorithm", 5.7, 2.2),
        ("Key Metric", 8.0, 2.0), ("Integration", 10.1, 2.6)]
for label, x, w in cols:
    add_text(slide, label, Inches(x), header_y + Inches(0.05), Inches(w), Inches(0.35),
             font_size=12, color=WHITE, bold=True)

for i, (num, name, approach, algo, metric, integration) in enumerate(pipelines):
    row_y = header_y + Inches(0.5) + Inches(i * 0.55)
    bg_color = CARD_BG if i % 2 == 0 else LIGHT_GRAY
    add_shape(slide, Inches(0.6), row_y, Inches(12.1), Inches(0.5), fill=bg_color)
    vals = [num, name, approach, algo, metric, integration]
    for j, (label, x, w) in enumerate(cols):
        add_text(slide, vals[j], Inches(x), row_y + Inches(0.07), Inches(w), Inches(0.35),
                 font_size=12, color=SLATE_NAVY if j == 0 else DARK_TEXT, bold=(j==0))

add_rounded_rect(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.3), fill=CARD_BG)
add_text(slide, "All 5 pipelines run nightly at 2:00 AM via automated scheduler. Models serialized with joblib. Predictions served through .NET API as CSV lookups. Every pipeline follows the full textbook lifecycle: Problem Framing > Data Prep > Exploration > Modeling > Feature Selection > Evaluation > Deployment.",
         Inches(1.1), Inches(5.4), Inches(11.2), Inches(1.0),
         font_size=14, color=BODY_TEXT)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — ARCHITECTURE (kept brief — just tech stack)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "ARCHITECTURE", "How It's Built")

# Two-column layout
add_text(slide, "Frontend", Inches(0.8), Inches(2.0), Inches(5), Inches(0.5),
         font_size=18, color=SKY_BLUE, bold=True)
fe_items = [
    "React 19 + TypeScript + Vite",
    "Tailwind CSS v4",
    "Zustand (client state) + React Query (server state)",
    "Recharts for visualizations",
    "react-hook-form + Zod for validation",
    "Deployed on Vercel",
]
add_bullet_list(slide, fe_items, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5),
                font_size=14, color=BODY_TEXT,
                bullet_color=(0x6A, 0xA8, 0xCB))

add_text(slide, "Backend & Data", Inches(7), Inches(2.0), Inches(5), Inches(0.5),
         font_size=18, color=SAGE_GREEN, bold=True)
be_items = [
    ".NET 10 / C# Web API + EF Core 9",
    "PostgreSQL on Supabase",
    "ASP.NET Identity + JWT authentication",
    "Python + scikit-learn + statsmodels for ML",
    "joblib model serialization + nightly refresh",
    "Cloud deployment (backend API)",
]
add_bullet_list(slide, be_items, Inches(7), Inches(2.5), Inches(5.5), Inches(3.5),
                font_size=14, color=BODY_TEXT,
                bullet_color=(0x5E, 0x8C, 0x5E))

add_slide_number(slide, 7, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — SECURITY & PRIVACY (expanded, no tech stack)
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "SECURITY & PRIVACY", "Protecting the Most Sensitive Data")

add_text(slide, "This system handles data about minors who are abuse survivors.\nPrivacy and data protection aren't optional -- they're foundational.",
         Inches(0.8), Inches(2.0), Inches(11), Inches(0.7),
         font_size=16, color=BODY_TEXT)

# Security cards - 3 columns
sec_col_w = Inches(3.6)
sec_gap = Inches(0.4)
sec_y = Inches(2.9)
sec_h = Inches(4.0)

# Column 1 - Authentication & Access
c1 = add_rounded_rect(slide, sx, sec_y, sec_col_w, sec_h, fill=CARD_BG)
add_shape(slide, sx, sec_y, sec_col_w, Inches(0.06), fill=SLATE_NAVY)
add_text(slide, "Authentication & Access", sx + Inches(0.3), sec_y + Inches(0.3),
         Inches(3), Inches(0.4), font_size=16, color=SLATE_NAVY, bold=True)
auth_items = [
    "ASP.NET Identity + JWT token auth",
    "Google OAuth third-party login",
    "Multi-Factor Authentication (TOTP app)",
    "Custom password policy (class spec)",
    "Role-Based Access Control (Admin / Donor / Public)",
    "All CUD endpoints require Admin role",
    "Delete confirmation on every destructive action",
]
add_bullet_list(slide, auth_items, sx + Inches(0.3), sec_y + Inches(0.8),
                Inches(3), Inches(3.0),
                font_size=12, color=BODY_TEXT, bullet_color=(0x2D, 0x3A, 0x4A), spacing=1.5)

# Column 2 - Privacy & Compliance
c2_x = sx + sec_col_w + sec_gap
c2 = add_rounded_rect(slide, c2_x, sec_y, sec_col_w, sec_h, fill=CARD_BG)
add_shape(slide, c2_x, sec_y, sec_col_w, Inches(0.06), fill=SAGE_GREEN)
add_text(slide, "Privacy & Compliance", c2_x + Inches(0.3), sec_y + Inches(0.3),
         Inches(3), Inches(0.4), font_size=16, color=SAGE_GREEN, bold=True)
privacy_items = [
    "GDPR-compliant privacy policy (customized)",
    "Functional GDPR cookie consent banner",
    "Browser-accessible cookie for user preferences",
    "Language and currency preferences saved",
    "Dark/light mode via cookie (nd_theme)",
    "Data sanitization on all inputs",
    "Credentials stored in env vars, not in repo",
]
add_bullet_list(slide, privacy_items, c2_x + Inches(0.3), sec_y + Inches(0.8),
                Inches(3), Inches(3.0),
                font_size=12, color=BODY_TEXT, bullet_color=(0x5E, 0x8C, 0x5E), spacing=1.5)

# Column 3 - Infrastructure Security
c3_x = c2_x + sec_col_w + sec_gap
c3 = add_rounded_rect(slide, c3_x, sec_y, sec_col_w, sec_h, fill=CARD_BG)
add_shape(slide, c3_x, sec_y, sec_col_w, Inches(0.06), fill=SKY_BLUE)
add_text(slide, "Infrastructure", c3_x + Inches(0.3), sec_y + Inches(0.3),
         Inches(3), Inches(0.4), font_size=16, color=SKY_BLUE, bold=True)
infra_items = [
    "HTTPS/TLS encryption on all connections",
    "HTTP to HTTPS redirect enforced",
    "HSTS (HTTP Strict Transport Security)",
    "Content-Security-Policy header via middleware",
    "Real DBMS (PostgreSQL) for identity store",
    "Deployed publicly on cloud infrastructure",
    "Sensitive endpoints require authentication",
]
add_bullet_list(slide, infra_items, c3_x + Inches(0.3), sec_y + Inches(0.8),
                Inches(3), Inches(3.0),
                font_size=12, color=BODY_TEXT, bullet_color=(0x6A, 0xA8, 0xCB), spacing=1.5)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — BUSINESS CASE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "THE CASE FOR INVESTMENT", "Why New Dawn Deserves Continued Investment")

stats = [
    ("5", "Production ML\nPipelines"),
    ("17", "Database Tables\nSeeded"),
    ("60+", "Residents\nManaged"),
    ("6", "Social Platforms\nSupported"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + Inches(i * 1.7)
    add_text(slide, num, x, Inches(2.3), Inches(1.5), Inches(0.7),
             font_size=36, color=SLATE_NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, label, x, Inches(3.0), Inches(1.5), Inches(0.8),
             font_size=12, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, "Value Delivered", Inches(7.2), Inches(2.1), Inches(5), Inches(0.5),
         font_size=18, color=SLATE_NAVY, bold=True)
value_items = [
    "Replaces manual risk tracking with ML-powered early warning",
    "Eliminates guesswork from social media with real-time predictions",
    "Predicts donor lapse before it happens for proactive retention",
    "Full case lifecycle management in one place for limited staff",
    "Language and currency preferences for international donors",
    "All pipelines refresh nightly -- the system gets smarter over time",
    "Enterprise-grade security for the most sensitive data",
]
add_bullet_list(slide, value_items, Inches(7.2), Inches(2.6), Inches(5.5), Inches(3.5),
                font_size=13, color=BODY_TEXT,
                bullet_color=(0x2D, 0x3A, 0x4A))

add_text(slide, "\"New Dawn doesn't just manage data -- it transforms data into decisions\nthat protect girls and grow the mission.\"",
         Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
         font_size=18, color=SKY_BLUE, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

add_slide_number(slide, 9, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING / Q&A
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, 0, 0, W, Inches(0.08), fill=ACCENT_BAR)

add_text(slide, "NEW DAWN", Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
         font_size=56, color=SLATE_NAVY, bold=True, alignment=PP_ALIGN.CENTER,
         font_name="Calibri Light")

add_accent_line(slide, Inches(5.5), Inches(3.1), Inches(2.3), GOLDEN_HONEY, Pt(3))

add_text(slide, "A Path to Healing and Hope", Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
         font_size=28, color=SKY_BLUE, alignment=PP_ALIGN.CENTER, font_name="Calibri Light")

add_text(slide, "Thank You  |  Questions?", Inches(1.5), Inches(4.6), Inches(10), Inches(0.8),
         font_size=24, color=SLATE_NAVY, alignment=PP_ALIGN.CENTER, bold=True)

add_text(slide, "new-dawn-virid.vercel.app", Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, "Section 3  |  Group 9\nZack Hada   |   Lincoln Lyons   |   [Team Members]",
         Inches(1.5), Inches(6.1), Inches(10), Inches(0.8),
         font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ── Save ────────────────────────────────────────────────────
output_path = "/Users/zackhada/Documents/coding/school/New-Dawn/presentation/New_Dawn_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
